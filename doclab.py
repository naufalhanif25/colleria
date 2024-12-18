# Importing necessary libraries and modules
import docx2pdf
import pdfplumber
from PyPDF2 import PdfReader
import docx
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from PIL import Image
import fitz
import img2pdf
import pandas as pd
import pptxtopdf
from pptx import Presentation
from pptx.util import Inches
import re
import io
import os

# Function to run the doclab algorithm (conversion process)
def doclab(input_path, output_path, input_ext, output_ext):
    """
    This function is used to convert any 
    type of document into another extension
    Functions:
    - docx_2_pdf (convert .docx to .pdf)
    - pdf_2_txt (convert .pdf to .txt)
    - pdf_2_docx (convert .pdf to .docx)
    - pdf_2_image (convert .pdf to images)
    - image_2_pdf (convert image to .pdf)
    """

    # Extract base file name and extension and modify the input path
    temp_input_path = os.path.basename(input_path)
    input_name, input_ext = os.path.splitext(temp_input_path)
    output_path = f"{output_path}/{input_name}{output_ext}"

    # Function to convert docx to pdf
    def docx_2_pdf(input_path, output_path):
        """
        This function is responsible for changing 
        the extension of a file from .docx to .pdf
        """

        # Convert docx to pdf
        docx2pdf.convert(input_path, output_path)

    # Function to convert pdf to docx
    def pdf_2_docx(input_path, output_path):
        """
        This function is responsible for changing 
        the extension of a file from .pdf to .docx
        """

        # Create a new DOCX document
        doc = docx.Document()

        # Read the PDF file
        with pdfplumber.open(input_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                # Extract text from the page
                text = page.extract_text()

                # Add text to the DOCX document
                if text:
                    for index, line in enumerate(text.split("\n")):
                        # fixed_line = re.sub(r"[^\x00-\x7F]+|\x0c", " ", line)

                        if index == 0:
                            doc.add_heading(line, level = 1)  # argument: line or fixed_line
                        else:
                            paragraph = doc.add_paragraph(line)  # argument: line or fixed_line

                            # Set alignment (example: center alignment) 
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

                # Extract and add images
                for image_index, image in enumerate(page.images):
                    try:
                        # Extract the image bbox and convert to PIL image
                        x0, y0, x1, y1 = image["x0"], image["top"], image["x1"], image["bottom"]
                        image_bbox = (x0, y0, x1, y1)
                        img = page.within_bbox(image_bbox).to_image(resolution = 300)

                        # Save the PIL image
                        image_path = f"{os.path.dirname(output_path)}/{input_name}_image_{page_num}_{image_index}.png"
                        img.save(image_path, format = "PNG", bits = 32)

                        # Open the image with Pillow and adjust the resolution 
                        with Image.open(image_path) as img: 
                            width, height = img.size

                            img = img.resize((width, height), Image.Resampling.LANCZOS) 
                            img.save(image_path, format = "PNG", bits = 32)

                            # Add the image to the DOCX document
                            paragraph = doc.add_paragraph()
                            run = paragraph.add_run()

                            run.add_picture(image_path, width = Inches(0.003 * width), height = Inches(0.003 * height))
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

                        os.remove(image_path)
                    except Exception as e:
                        print(f"Error extracting image on page {page_num}: {e}")

                # Extract and add tables
                for table in page.extract_tables():
                    if table:
                        table_obj = doc.add_table(rows = len(table), cols = len(table[0]))

                        for row_index, row in enumerate(table):
                            for col_index, cell in enumerate(row):
                                table_obj.cell(row_index, col_index).text = cell

        # Save the DOCX document
        doc.save(output_path)
    
    # Function to convert pdf to pptx
    def pdf_2_pptx(input_path, output_path):
        """
        This function is responsible for changing 
        the extension of a file from .pdf to .pptx
        """

        # Load the PDF file
        with pdfplumber.open(input_path) as pdf:
            # Create a PowerPoint presentation
            presentation = Presentation()
            
            for page_num, page in enumerate(pdf.pages):
                # Add a slide for each page
                slide_layout = presentation.slide_layouts[5]  # Blank layout
                slide = presentation.slides.add_slide(slide_layout)
                
                # Extract text from the page
                text = page.extract_text()
                if text:
                    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
                    text_frame = textbox.text_frame

                    for line in text.split("\n"):
                        paragraph = text_frame.add_paragraph()
                        paragraph.text = line

                # Extract images from the page
                for img_num, img in enumerate(page.images):
                    x0, y0, x1, y1 = img["x0"], img["top"], img["x1"], img["bottom"]
                    img_bbox = (x0, y0, x1, y1)
                    extracted_img = page.within_bbox(img_bbox).to_image(resolution = 300)

                    # Save the image to a byte stream
                    image_stream = io.BytesIO()
                    extracted_img.save(image_stream, format = "PNG", bits = 32)
                    image_stream.seek(0)

                    # Add image to slide
                    slide.shapes.add_picture(image_stream, Inches(1), Inches(3), width = Inches(3))

            # Iterate through each slide in the presentation
            for slide in presentation.slides:
                for shape in slide.shapes:  # Iterate through each shape in the current slide
                    for placeholder in slide.shapes.placeholders:  # Iterate through each placeholder in the current slide
                        # Get the underlying XML element of the placeholder
                        slide_placeholder = placeholder._sp

                        # Remove the XML element from its parent
                        slide_placeholder.getparent().remove(slide_placeholder)


        # Save the PowerPoint presentation
        presentation.save(output_path)
    
    # Function to convert pdf to image
    def pdf_2_image(input_path, output_path):
        """
        This function is responsible for changing 
        the extension of a file from .pdf to images
        (.png, .jpg, .jpeg, .tiff)
        """

        # Extract base file name and extension and modify the output path
        temp_output_path = os.path.basename(output_path)
        output_name, output_ext = os.path.splitext(temp_output_path)

        # Open the PDF file
        pdf = fitz.open(input_path)

        for page_num in range(len(pdf)):
            # Get a page
            page = pdf[page_num]

            # Render page to an image (300 dpi)
            pix = page.get_pixmap(dpi = 300)

            # Save the image
            pix.save(f"{os.path.dirname(output_path)}/{output_name}_{page_num + 1}{output_ext}")

    # Function to convert image to pdf
    def image_2_pdf(input_path, output_path):
        """
        This function is responsible for changing 
        the extension of a file from images
        (.png, .jpg, .jpeg, .tiff) to .pdf
        """

        # opening image
        image = Image.open(input_path)

        # converting into chunks using img2pdf
        pdf_chunks = img2pdf.convert(image.filename)

        with open(output_path, "wb") as pdf:
            pdf.write(pdf_chunks)
            pdf.close()

        image.close()

    # Function to convert images to other image formats 
    def image_2_image(input_path, output_path):
        """
        This function is responsible for changing 
        the extension of a file from image
        (.png, .jpg, .jpeg, .tiff) to other 
        image formats (.png, .jpg, .jpeg, .tiff)
        """

        # Open the input image
        image = Image.open(input_path)

        # Convert RGBA to RGB if needed 
        if image.mode == "RGBA": 
            image = image.convert("RGB")
        
        # Save the image in the specified output format
        image.save(output_path)

    # Function to convert xlsx to csv
    def xlsx_2_csv(input_path, output_path):
        """
        This function is responsible for changing 
        the extension of a file from .xlsx to .csv
        """

        # Read the Excel file
        data = pd.read_excel(input_path)

        # Save to CSV file
        data.to_csv(output_path, index = False)

    # Function to convert csv to xlsx
    def csv_2_xlsx(input_path, output_path):
        """
        This function is responsible for changing 
        the extension of a file from .csv to .xlsx
        """

        # Read the Excel file
        data = pd.read_csv(input_path)

        # Save to CSV file
        data.to_excel(output_path, index = False)

    # Function to convert pptx to pdf
    def pptx_2_pdf(input_path, output_path):
        """
        This function is responsible for changing 
        the extension of a file from .pptx to .pdf
        """

        # Get path directory
        output_path = f"{os.path.dirname(output_path)}"

        # Convert pptx to pdf
        pptxtopdf.convert(input_path, output_path)

    # List of supported file extensions
    IMG_EXT = [".png", ".jpg", ".jpeg"]
    EXT_DICT = {
        (".pdf", ".docx") : pdf_2_docx,
        (".pdf", ".doc") : pdf_2_docx,
        (".docx", ".pdf") : docx_2_pdf,
        (".doc", ".pdf") : docx_2_pdf,
        (".xlsx", ".csv") : xlsx_2_csv,
        (".xls", ".csv") : xlsx_2_csv,
        (".csv", ".xlsx") : csv_2_xlsx,
        (".csv", ".xls") : csv_2_xlsx,
        (".pptx", ".pdf") : pptx_2_pdf,
        (".ppt", ".pdf") : pptx_2_pdf,
        (".pdf", ".pptx") : pdf_2_pptx,
        (".pdf", ".ppt") : pdf_2_pptx
    }

    # Adding image extensions dynamically to the list
    for ext in IMG_EXT: 
        EXT_DICT[(ext, ".pdf")] = image_2_pdf 
        EXT_DICT[(".pdf", ext)] = pdf_2_image

    for ext_left in IMG_EXT:
        for ext_right in IMG_EXT:
            if ext_left != ext_right and (ext_left, ext_right) not in EXT_DICT:
                EXT_DICT[(ext_left, ext_right)] = image_2_image

    if (input_ext, output_ext) in EXT_DICT:
        EXT_DICT[(input_ext, output_ext)](input_path, output_path)
